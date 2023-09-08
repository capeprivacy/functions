from io import BufferedReader
from io import StringIO
import logging
import mimetypes
import os
from typing import Optional, Tuple
import uuid

import docx2txt
import fitz
import pandas
import pptx
import json
import base64

logger = logging.getLogger("reader")

excel_mimetypes = [
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel.sheet.macroEnabled.12",
    "application/vnd.ms-excel.sheet.binary.macroEnabled.12",
]

mimetypes_to_human = {
    "application/vnd.ms-excel": "spreadsheet",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "spreadsheet",
    "application/vnd.ms-excel.sheet.macroEnabled.12": "spreadsheet",
    "application/vnd.ms-excel.sheet.binary.macroEnabled.12": "spreadsheet",
    "application/octet-stream": "bytes",
    "text/markdown": "markdown",
    "text/plain": "plaintext",
    "text/csv": "csv",
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "powerpoint",
}


class UnsupportedFileError(Exception):
    pass


def extract_text_from_filepath(
    filepath: str, mimetype: Optional[str] = None
) -> Tuple[str, str]:
    """Return the text content of a file given its filepath."""

    if not mimetype or mimetype == "application/octet-stream":
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)
        logger.info(f"guessed mimetype {mimetype}")

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise UnsupportedFileError(
                "Unsupported file type, could not determine mimetype"
            )

    try:
        with open(filepath, "rb") as file:
            # extracting text can be slow so run in threadpool
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(f"Error extracting text from file: {e}")
        raise e

    return extracted_text, mimetype


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyMuPDF
        with fitz.open(file) as pdf:
            extracted_text = " ".join([page.get_text() for page in pdf])

    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        try:
            # if there are any errors with decoding then replace those bytes with `?`.
            # this handles the situation where a file may contain invalid bytes but is otherwise
            # correct
            extracted_text = file.read().decode("utf-8", errors="replace")
        except UnicodeDecodeError as e:
            logger.error(f"uploaded file was not utf-8 encoded {e}")
            raise UnsupportedFileError("uploaded file was not utf-8 encoded")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""

        # if there are any errors with decoding then replace those bytes with `?`.
        # this handles the situation where a file may contain invalid bytes but is otherwise
        # correct
        for line in file:
            extracted_text += line.decode("utf-8", errors="replace")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    elif mimetype in excel_mimetypes:
        extracted_text = ""
        df = pandas.read_excel(file, sheet_name=None)
        for j in df:
            b = StringIO()
            df[j].to_csv(b)
            extracted_text += b.getvalue() + "\n\n"
    else:
        # Unsupported file type
        raise UnsupportedFileError(f"Unsupported file type: {mimetype}")

    return extracted_text


# Extract text from a file based on its mimetype
def extract_text_from_form_file(data, mimetype, filename):
    _, ext = os.path.splitext(filename)
    temp_file_path = f"temp_file-{uuid.uuid4()}{ext}"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(data)

    try:
        extracted_text, mimetype = extract_text_from_filepath(
            temp_file_path, mimetype
        )
    except Exception as e:
        logger.error(f"Error extracting text from filepath: {e}")
        os.remove(temp_file_path)

        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text


def cape_handler(data):
    inp = json.loads(data)
    d = base64.b64decode(inp["data"])
    return extract_text_from_form_file(d, inp["mimetype"], inp["filename"])
